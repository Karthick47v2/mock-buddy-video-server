"""Extract details from slides"""

from urllib.error import HTTPError
from urllib import request
from collections import Counter

from pptx import Presentation


class Slide:
    """Class for analyzing presentation slide"""

    def __init__(self, url, name):
        """Initialize slide details

        Args:
            url (str): shared url
            name (str): name for slide
        """
        self.__url = url.split('/')
        self.__url[-1] = 'export/pptx'
        self.__url = '/'.join(self.__url)
        self.filename = name

        self.__txt_lst = []
        self.__word_count = []
        self.__shape_count = 0
        self.__slide_count = 0
        self.__font_lst = []

        self.__load_pptx()

    def __load_pptx(self):
        """Download slide from source

        Returns:
            dict[str, str]: error status if anything happens
        """
        try:
            request.urlretrieve(self.__url, self.filename)
            return {
                'status': 200
            }
        except HTTPError:
            return {
                'status': 400
            }

    def __set_slide_count(self, slides):
        """Update slide count

        Args:
            slides (obj): pptx slides object
        """
        self.__slide_count = len(slides)

    def __extract_txt(self, runs, word_count, txt):
        """Extract text from each slide

        Args:
            runs (obj): paragraph object
            word_count (list[int]): list of word count per para
            txt (list[str]): list of text per para
        """
        for run in runs:
            if run.font.name is not None:
                self.__font_lst.append(run.font.name)

            word_count.append(
                len(run.text.split(' ')))
            txt.append(run.text)

    def __increment_shape_count(self, shape_type):
        """Update shape count

        Args:
            shape_type (int): type of shape (picture = 13, auto-shape = 6, etc.)
        """
        if shape_type in [13, 6]:
            self.__shape_count += 1

    def extract_pptx(self):
        """Extract required details from slide (text, shape count, word count, etc.)"""
        with open(self.filename, 'rb') as pptx:
            prs = Presentation(pptx)
            self.__set_slide_count(prs.slides)

            for slide in prs.slides:
                slide_txt_lst = []
                slide_word_count = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            self.__extract_txt(
                                para.runs, slide_word_count, slide_txt_lst)
                    self.__increment_shape_count(shape.shape_type)

                self.__txt_lst.append(slide_txt_lst)
                self.__word_count.append(slide_word_count)

    def get_txt(self):
        """Return all text from pptx

        Returns:
            list(list(str)): list of text for each slide
        """
        return self.__txt_lst

    def get_data(self):
        """Return infomation about slide other than text.

        Returns:
            Tuple(int, list(list(int)), int, dict(str: int)): tuple of number of slides,
            word count per slide, shape count, unique value count of fonts
        """
        return (self.__slide_count,
                self.__word_count,
                self.__shape_count,
                dict(Counter(self.__font_lst)))
